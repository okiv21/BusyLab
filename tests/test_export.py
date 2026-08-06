"""Export tests.

Spec Pillar 6 wants static, portable, universally openable output, and
explicitly rejects rendered MP4. So the checks here are that the files are
genuinely valid rather than merely non-empty, that a deck carries native
editable charts rather than pictures, and that a finding whose shape has no
honest chart gets a table instead of a misleading one.
"""

from __future__ import annotations

import io

import pytest
from fastapi.testclient import TestClient

from busylab.analysis import analyse
from busylab.export import MAX_SLIDES, series_for, table_for, to_pdf, to_pptx
from busylab.findings import ChartType, Evidence, Finding, FindingType

from . import fixtures


@pytest.fixture(scope="module")
def story():
    return analyse(fixtures.planted_business(), strict=True)


# --------------------------------------------------------------------------
# The files are real
# --------------------------------------------------------------------------


def test_pdf_is_a_valid_document(story) -> None:
    payload = to_pdf(story.findings, business_name="Fern & Flame")

    assert payload.startswith(b"%PDF-")
    assert payload.rstrip().endswith(b"%%EOF")
    assert len(payload) > 3000


def test_pptx_is_a_valid_openable_deck(story) -> None:
    from pptx import Presentation

    payload = to_pptx(story.findings, business_name="Fern & Flame")
    assert payload.startswith(b"PK\x03\x04"), "a pptx is a zip archive"

    deck = Presentation(io.BytesIO(payload))
    # A cover plus one slide per finding.
    assert len(deck.slides) == min(len(story.findings), MAX_SLIDES) + 1


def test_the_deck_carries_native_editable_charts(story) -> None:
    """A real chart object can be restyled in PowerPoint; a picture cannot."""
    from pptx import Presentation

    deck = Presentation(io.BytesIO(to_pptx(story.findings)))
    charts = sum(1 for s in deck.slides for shape in s.shapes if shape.has_chart)

    assert charts >= 4, "most findings should map onto a native chart"


def test_findings_without_an_honest_chart_get_a_table(story) -> None:
    """A cohort triangle has no bar-chart form. A table beats a bad picture."""
    from pptx import Presentation

    deck = Presentation(io.BytesIO(to_pptx(story.findings)))
    tables = sum(1 for s in deck.slides for shape in s.shapes if shape.has_table)

    assert tables >= 1


def test_every_finding_reaches_the_deck_one_way_or_another(story) -> None:
    from pptx import Presentation

    deck = Presentation(io.BytesIO(to_pptx(story.findings)))
    # Every slide past the cover must carry the finding's sentence.
    for slide in list(deck.slides)[1:]:
        text = " ".join(
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert len(text.strip()) > 20


def test_a_long_story_is_trimmed_to_a_readable_length() -> None:
    """A deck nobody finishes is a deck nobody opens twice."""
    from pptx import Presentation

    many = [
        Finding(
            id=f"f{i}",
            type=FindingType.RANKING,
            summary=f"Finding number {i} says something factual.",
            facts={"item_count": 3},
            evidence=Evidence(method="test"),
            chart_data={"bars": [{"label": "a", "value": 1.0}]},
        )
        for i in range(30)
    ]
    deck = Presentation(io.BytesIO(to_pptx(many)))
    assert len(deck.slides) == MAX_SLIDES + 1


def test_an_empty_story_still_produces_a_valid_file() -> None:
    assert to_pdf([]).startswith(b"%PDF-")
    assert to_pptx([]).startswith(b"PK\x03\x04")


# --------------------------------------------------------------------------
# Chart mapping
# --------------------------------------------------------------------------


def _finding(chart: ChartType, chart_data: dict, ftype=FindingType.RANKING) -> Finding:
    return Finding(
        id="x",
        type=ftype,
        summary="A sentence.",
        chart=chart,
        chart_data=chart_data,
        evidence=Evidence(method="test"),
    )


def test_bars_map_to_a_bar_series() -> None:
    series = series_for(
        _finding(
            ChartType.BAR_HORIZONTAL,
            {"bars": [{"label": "A", "value": 10.0}, {"label": "B", "value": 5.0}]},
        )
    )
    assert series is not None and series.kind == "bar"
    assert series.values == [10.0, 5.0]


def test_a_donut_maps_to_a_pie_series() -> None:
    series = series_for(
        _finding(ChartType.DONUT, {"slices": [{"label": "A", "value": 7.0}]})
    )
    assert series is not None and series.kind == "pie"


def test_a_trend_maps_to_a_line_series() -> None:
    series = series_for(
        _finding(
            ChartType.LINE_WITH_BAND,
            {"series": [{"period": "2025-01-01", "value": 3.0}]},
        )
    )
    assert series is not None and series.kind == "line"


def test_a_forecast_line_includes_the_projection() -> None:
    series = series_for(
        _finding(
            ChartType.FORECAST_FAN,
            {
                "history": [{"period": "2025-01-01", "value": 10.0}],
                "forecast": [{"period": "2025-02-01", "mean": 12.0}],
            },
        )
    )
    assert series is not None
    assert series.values == [10.0, 12.0]


def test_a_cohort_heatmap_has_no_chart_form() -> None:
    assert series_for(_finding(ChartType.COHORT_HEATMAP, {"curve": []})) is None


def test_a_quadrant_has_no_chart_form() -> None:
    assert series_for(_finding(ChartType.QUADRANT, {"customers": []})) is None


def test_a_table_carries_readable_labels_and_values(story) -> None:
    finding = next(f for f in story.findings if f.id == "rfm_segments")
    rows = table_for(finding)

    assert rows
    for label, value in rows:
        assert label[0].isupper(), "labels are humanised, not raw keys"
        assert "_" not in label
        assert value


def test_percentages_are_rendered_as_percentages(story) -> None:
    finding = next(f for f in story.findings if f.id == "concentration")
    rows = dict(table_for(finding))
    percentage_rows = [v for k, v in rows.items() if "share" in k.lower()]
    assert any(v.endswith("%") for v in percentage_rows)


# --------------------------------------------------------------------------
# Over HTTP
# --------------------------------------------------------------------------


@pytest.fixture
def client(tmp_path, monkeypatch):
    monkeypatch.setenv("BUSYLAB_INLINE_WORKER", "0")
    monkeypatch.delenv("GROQ_API_KEY", raising=False)

    from api import main
    from api.handlers import build_handlers
    from api.jobs import JobStore, Worker
    from api.storage import LocalFileStore

    files = LocalFileStore(root=tmp_path / "storage")
    store = JobStore(tmp_path / "export.db")
    worker = Worker(store, build_handlers(files))
    monkeypatch.setattr(main, "_store", store)
    monkeypatch.setattr(main, "_worker", worker)
    monkeypatch.setattr(main, "FILES", files)

    with TestClient(main.app) as test_client:
        test_client.worker = worker
        test_client.store = store
        yield test_client


@pytest.fixture
def analysed(client, tmp_path):
    path = tmp_path / "fern_and_flame.xlsx"
    fixtures.planted_business().drop(columns=["salesperson"]).to_excel(
        path, index=False
    )
    with open(path, "rb") as handle:
        upload = client.post(
            "/uploads", files={"file": (path.name, handle, "application/vnd.ms-excel")}
        ).json()
    client.worker.drain()
    client.post(f"/datasets/{upload['dataset_id']}/columns", json={"roles": {}})
    client.worker.drain()
    return upload["dataset_id"]


def test_pdf_downloads_with_a_sensible_filename(client, analysed) -> None:
    response = client.get(f"/datasets/{analysed}/export.pdf")

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"
    assert "fern_and_flame" in response.headers["content-disposition"]
    assert response.content.startswith(b"%PDF-")


def test_pptx_downloads_with_the_right_content_type(client, analysed) -> None:
    response = client.get(f"/datasets/{analysed}/export.pptx")

    assert response.status_code == 200
    assert "presentation" in response.headers["content-type"]
    assert response.content.startswith(b"PK\x03\x04")


def test_an_unknown_format_is_refused(client, analysed) -> None:
    response = client.get(f"/datasets/{analysed}/export.mp4")
    assert response.status_code == 404
    assert "pdf or pptx" in response.json()["detail"]


def test_exporting_before_analysis_is_a_conflict(client, tmp_path) -> None:
    path = tmp_path / "pending.xlsx"
    fixtures.base_sales(n=200).to_excel(path, index=False)
    with open(path, "rb") as handle:
        upload = client.post(
            "/uploads", files={"file": (path.name, handle, "application/vnd.ms-excel")}
        ).json()
    client.worker.drain()

    assert client.get(f"/datasets/{upload['dataset_id']}/export.pdf").status_code == 409


def test_a_held_analysis_cannot_be_exported(client, tmp_path) -> None:
    """Exporting a story the quality gate refused to publish would launder it."""
    import pandas as pd

    frame = fixtures.planted_business().drop(columns=["salesperson"])
    doubled = pd.concat([frame, frame], ignore_index=True)

    path = tmp_path / "doubled.xlsx"
    doubled.to_excel(path, index=False)
    with open(path, "rb") as handle:
        upload = client.post(
            "/uploads", files={"file": (path.name, handle, "application/vnd.ms-excel")}
        ).json()
    client.worker.drain()
    client.post(f"/datasets/{upload['dataset_id']}/columns", json={"roles": {}})
    client.worker.drain()

    response = client.get(f"/datasets/{upload['dataset_id']}/export.pdf")
    assert response.status_code == 409
    assert "quality gate" in response.json()["detail"]


def test_exports_on_a_missing_dataset_are_a_404(client) -> None:
    assert client.get("/datasets/nope/export.pdf").status_code == 404
