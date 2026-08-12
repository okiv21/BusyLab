"""Pillar 6: static, portable outputs for people who are not in the app.

A board or an investor gets a PDF or a deck, not a login. Both are cheap to
generate, universally openable, and carry the same ranked story the app shows.

**Rendered MP4 video is explicitly rejected** (spec Pillar 6): slow to generate,
expensive per business and per refresh, and frozen the moment the data changes.
It fights the live-and-current principle that the rest of the product is built
on. Present mode covers the same need without any of that, because it plays
inside the app and updates the instant the data does.

Two deliberate choices here:

**Native charts in the deck, not pictures.** python-pptx can emit a real
PowerPoint chart object, which stays editable — someone can restyle it for
their own template or fix a label without coming back to us. A screenshot
cannot. Where a finding's shape does not map onto a native chart type (a cohort
triangle, a correlation matrix) it gets a clean table instead of a bad picture.

**Pure Python, no native dependencies.** reportlab and python-pptx are both
pip-installable with nothing to compile, which matters because memory and
install size are the deployment constraint (spec 9).
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from datetime import date
from typing import Any, Iterable

from .findings import ChartType, Finding, Severity

# Brand palette, matching the app so an export does not look like a different
# product.
INK = (0x21, 0x1C, 0x15)
INK_MUTED = (0x6E, 0x67, 0x5C)
INK_LIGHT = (0x8A, 0x83, 0x78)
ACCENT = (0xE8, 0x5A, 0x32)
ACCENT_DARK = (0xC7, 0x47, 0x22)
GOOD = (0x1F, 0xA9, 0x7A)
WARN = (0xB0, 0x6A, 0x1E)
PAGE = (0xFB, 0xFA, 0xF8)
LINE = (0xF0, 0xEB, 0xE3)

#: Chart shapes that map onto a native slide chart.
_NATIVE_BAR = {ChartType.BAR_HORIZONTAL, ChartType.GROUPED_BARS}
_NATIVE_LINE = {ChartType.LINE_WITH_BAND, ChartType.FORECAST_FAN, ChartType.STACKED_AREA}
_NATIVE_PIE = {ChartType.DONUT, ChartType.TREEMAP}

#: How many findings a board deck should carry before it stops being read.
MAX_SLIDES = 12


@dataclass
class Series:
    """A chart reduced to labels and values, which is all an export needs."""

    kind: str  # "bar", "line" or "pie"
    labels: list[str]
    values: list[float]
    title: str = ""
    #: True when the values are signed contributions either side of zero, so up
    #: and down have to read as opposite rather than as one colour of bar. A
    #: ranking is not diverging: every bar is positive and colouring them by
    #: value would double-encode length as hue.
    diverging: bool = False


def _severity_colour(severity: Severity) -> tuple[int, int, int]:
    return {
        Severity.URGENT: ACCENT_DARK,
        Severity.WATCH: WARN,
        Severity.GOOD: GOOD,
        Severity.NEUTRAL: INK_MUTED,
    }.get(severity, INK_MUTED)


def _compact(value: float) -> str:
    magnitude = abs(value)
    if magnitude >= 1_000_000_000:
        return f"{value / 1_000_000_000:,.1f}b"
    if magnitude >= 1_000_000:
        return f"{value / 1_000_000:,.1f}m"
    if magnitude >= 1_000:
        return f"{value / 1_000:,.1f}k"
    return f"{value:,.0f}"


def series_for(finding: Finding) -> Series | None:
    """Reduce a finding's chart payload to something an export can draw.

    Returns None where the shape genuinely does not fit a bar, line or pie.
    Those findings get a table, which is honest, rather than a chart that
    misrepresents them.
    """
    data = finding.chart_data or {}
    chart = finding.chart

    if chart in _NATIVE_BAR:
        bars = data.get("bars") or data.get("groups") or []
        if bars:
            return Series(
                "bar",
                [str(b.get("label", "")) for b in bars][:8],
                [float(b.get("value", 0)) for b in bars][:8],
            )

    if chart in _NATIVE_PIE:
        slices = data.get("slices") or []
        if slices:
            return Series(
                "pie",
                [str(s.get("label", "")) for s in slices][:6],
                [float(s.get("value", 0)) for s in slices][:6],
            )

    if chart in _NATIVE_LINE:
        series = data.get("series") or []
        if series and "value" in (series[0] if series else {}):
            return Series(
                "line",
                [str(p.get("period", "")) for p in series],
                [float(p.get("value", 0)) for p in series],
            )
        # A forecast: the actual history, then the projected central line.
        history = data.get("history") or []
        forecast = data.get("forecast") or []
        if history:
            labels = [str(p.get("period", "")) for p in history]
            values = [float(p.get("value", 0)) for p in history]
            labels += [str(p.get("period", "")) for p in forecast]
            values += [float(p.get("mean", 0)) for p in forecast]
            return Series("line", labels, values)
        # Repeat versus new arrives as two stacked series; total is enough
        # for a static page, and the sentence carries the split.
        if series and "repeat" in (series[0] if series else {}):
            return Series(
                "line",
                [str(p.get("period", "")) for p in series],
                [
                    float(p.get("repeat", 0)) + float(p.get("new", 0))
                    for p in series
                ],
            )

    if chart is ChartType.WATERFALL:
        steps = [s for s in (data.get("steps") or []) if float(s.get("change", 0))]
        if steps:
            # Sorted by size before truncating, which the export was not doing.
            # The payload arrives sorted ascending by change, so taking the
            # first eight kept the five largest falls and dropped the two
            # largest rises entirely: on a real file that turned "one product
            # fell further than the whole change while others rose" into a page
            # where every bar pointed the same way. The chart disagreed with the
            # sentence printed directly above it.
            steps = sorted(
                steps, key=lambda s: abs(float(s.get("change", 0))), reverse=True
            )[:7]
            return Series(
                "bar",
                [str(s.get("label", "")) for s in steps],
                [float(s.get("change", 0)) for s in steps],
                diverging=True,
            )

    return None


def table_for(finding: Finding) -> list[tuple[str, str]]:
    """Key figures for a finding whose shape has no chart in an export."""
    facts = finding.facts or {}
    skip = {
        "contributions",
        "ranking",
        "means",
        "monthly_index",
        "pairs",
        "curve",
        "segments",
        "customers",
        "cohorts",
        "pace",
        "gap_drivers",
        "goal",
        "strongest",
    }
    rows: list[tuple[str, str]] = []
    for key, value in facts.items():
        if key in skip or isinstance(value, (list, dict)) or value is None:
            continue
        label = key.replace("_", " ").capitalize()
        if isinstance(value, bool):
            rows.append((label, "yes" if value else "no"))
        elif isinstance(value, float):
            if key.endswith("_pct") or key.endswith("_share") or "share" in key:
                rows.append((label, f"{value * 100:.0f}%"))
            else:
                rows.append((label, _compact(value)))
        elif isinstance(value, int):
            rows.append((label, f"{value:,}"))
        else:
            rows.append((label, str(value)[:60]))
        if len(rows) >= 6:
            break
    return rows


def _slides(findings: Iterable[Finding]) -> list[Finding]:
    """Trim to what a reader will actually get through."""
    return [f for f in findings][:MAX_SLIDES]


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------


def to_pdf(
    findings: list[Finding],
    *,
    business_name: str = "Your business",
    generated: date | None = None,
) -> bytes:
    """A portable report: one finding per block, chart where it fits."""
    from reportlab.lib.colors import Color
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.pdfgen import canvas as pdf_canvas

    def rgb(triple: tuple[int, int, int], alpha: float = 1.0) -> Color:
        r, g, b = (v / 255 for v in triple)
        return Color(r, g, b, alpha=alpha)

    buffer = io.BytesIO()
    page_width, page_height = A4
    pdf = pdf_canvas.Canvas(buffer, pagesize=A4)
    pdf.setTitle(f"{business_name} - business in review")

    margin = 18 * mm
    when = (generated or date.today()).isoformat()

    # Cover
    pdf.setFillColor(rgb(PAGE))
    pdf.rect(0, 0, page_width, page_height, fill=1, stroke=0)
    pdf.setFillColor(rgb(ACCENT))
    pdf.roundRect(margin, page_height - margin - 44, 34, 34, 10, fill=1, stroke=0)
    pdf.setFillColor(rgb((255, 255, 255)))
    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(margin + 10, page_height - margin - 34, "B")

    pdf.setFillColor(rgb(INK))
    pdf.setFont("Helvetica-Bold", 26)
    pdf.drawString(margin, page_height - margin - 96, "Business in review")
    pdf.setFont("Helvetica", 13)
    pdf.setFillColor(rgb(INK_MUTED))
    pdf.drawString(margin, page_height - margin - 118, f"{business_name} · {when}")

    pdf.setFont("Helvetica", 10.5)
    pdf.setFillColor(rgb(INK_LIGHT))
    _wrap(
        pdf,
        "Findings are ranked by how much they matter. Nothing here is advice: "
        "these are the facts as the numbers have them, and the decisions stay "
        "yours.",
        margin,
        page_height - margin - 150,
        page_width - 2 * margin,
        13,
    )

    y = page_height - margin - 200
    for index, finding in enumerate(_slides(findings), start=1):
        block = 150
        if y - block < margin:
            pdf.showPage()
            pdf.setFillColor(rgb(PAGE))
            pdf.rect(0, 0, page_width, page_height, fill=1, stroke=0)
            y = page_height - margin

        pdf.setFillColor(rgb(_severity_colour(finding.severity)))
        pdf.setFont("Helvetica-Bold", 9)
        pdf.drawString(margin, y, f"FINDING {index}".upper())

        pdf.setFillColor(rgb(INK))
        pdf.setFont("Helvetica-Bold", 12.5)
        y = _wrap(
            pdf, finding.summary, margin, y - 18, page_width - 2 * margin, 15
        )

        series = series_for(finding)
        if series and series.values:
            y = _draw_chart(pdf, series, margin, y - 10, page_width - 2 * margin, 74, rgb)
        else:
            rows = table_for(finding)
            pdf.setFont("Helvetica", 9.5)
            for label, value in rows[:4]:
                pdf.setFillColor(rgb(INK_LIGHT))
                pdf.drawString(margin + 4, y - 12, label)
                pdf.setFillColor(rgb(INK))
                pdf.drawRightString(margin + 220, y - 12, value)
                y -= 13

        pdf.setFont("Helvetica-Oblique", 8.5)
        pdf.setFillColor(rgb(INK_LIGHT))
        pdf.drawString(
            margin, y - 14, f"{finding.evidence.method} · {finding.evidence.strength}"
        )
        y -= 34

        pdf.setStrokeColor(rgb(LINE))
        pdf.setLineWidth(0.6)
        pdf.line(margin, y + 8, page_width - margin, y + 8)

    pdf.save()
    return buffer.getvalue()


def _wrap(pdf, text: str, x: float, y: float, width: float, leading: float) -> float:
    """Draw wrapped text, returning the y after the last line."""
    from reportlab.pdfbase.pdfmetrics import stringWidth

    font = pdf._fontname
    size = pdf._fontsize
    words = text.split()
    line = ""
    for word in words:
        candidate = f"{line} {word}".strip()
        if stringWidth(candidate, font, size) > width and line:
            pdf.drawString(x, y, line)
            y -= leading
            line = word
        else:
            line = candidate
    if line:
        pdf.drawString(x, y, line)
        y -= leading
    return y


def _draw_chart(pdf, series: Series, x: float, y: float, width: float, height: float, rgb) -> float:
    """Draw a small bar or line chart with reportlab primitives."""
    values = series.values
    if not values:
        return y

    top = max(values + [0.0])
    bottom = min(values + [0.0])
    span = (top - bottom) or 1.0
    baseline = y - height + (0 - bottom) / span * height

    if series.kind == "line":
        pdf.setStrokeColor(rgb(ACCENT))
        pdf.setLineWidth(1.6)
        step = width / max(len(values) - 1, 1)
        path = pdf.beginPath()
        for i, value in enumerate(values):
            px = x + i * step
            py = y - height + (value - bottom) / span * height
            if i == 0:
                path.moveTo(px, py)
            else:
                path.lineTo(px, py)
        pdf.drawPath(path)
    else:
        count = min(len(values), 8)
        slot = width / max(count, 1)
        bar_width = slot * 0.55
        for i in range(count):
            value = values[i]
            py = y - height + (min(value, 0) - bottom) / span * height
            bar_height = abs(value) / span * height
            # Green up against orange down only where the values are signed.
            # The previous pair was accent against warn, both warm and the
            # falling one duller, so a rise and a fall read as the same bar in
            # two shades. On a ranking, where every value is positive, one
            # colour is correct: colouring by value would encode bar length
            # twice and say nothing new.
            if series.diverging:
                fill = GOOD if value >= 0 else ACCENT
            else:
                fill = ACCENT
            centre = x + i * slot + slot / 2
            pdf.setFillColor(rgb(fill, 0.9))
            pdf.rect(
                x + i * slot + (slot - bar_width) / 2,
                py,
                bar_width,
                max(bar_height, 0.8),
                fill=1,
                stroke=0,
            )

            if series.diverging:
                # The figure, because a small bar cannot be read off the axis
                # and the reader has no tooltip on paper.
                pdf.setFillColor(rgb(ACCENT_DARK if value < 0 else GOOD))
                pdf.setFont("Helvetica-Bold", 6.5)
                above = value >= 0
                pdf.drawCentredString(
                    centre,
                    (py + bar_height + 3) if above else (py - 8),
                    # A plain hyphen, not a typographic minus. The base PDF
                    # fonts encode WinAnsi, which has no U+2212, so reportlab
                    # splits the string and draws the sign from a fallback font
                    # as a separately positioned run. On paper the difference is
                    # invisible; the fragility is not worth it.
                    f"{'+' if above else '-'}{_compact(abs(value))}",
                )

            pdf.setFillColor(rgb(INK_LIGHT))
            pdf.setFont("Helvetica", 6.5)
            label = series.labels[i] if i < len(series.labels) else ""
            pdf.drawCentredString(centre, y - height - 9, label[:14])

    pdf.setStrokeColor(rgb(LINE))
    pdf.setLineWidth(0.6)
    pdf.line(x, baseline, x + width, baseline)
    return y - height - 14


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------


def to_pptx(
    findings: list[Finding],
    *,
    business_name: str = "Your business",
    generated: date | None = None,
) -> bytes:
    """A slide deck with native, editable charts."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    def colour(triple: tuple[int, int, int]) -> RGBColor:
        return RGBColor(*triple)

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]
    when = (generated or date.today()).isoformat()

    def background(slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = colour(PAGE)

    # Cover
    cover = deck.slides.add_slide(blank)
    background(cover)
    box = cover.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.4))
    frame = box.text_frame
    frame.word_wrap = True
    title = frame.paragraphs[0]
    title.text = "Business in review"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = colour(INK)
    subtitle = frame.add_paragraph()
    subtitle.text = f"{business_name} · {when}"
    subtitle.font.size = Pt(18)
    subtitle.font.color.rgb = colour(INK_MUTED)
    note = frame.add_paragraph()
    note.text = (
        "Findings ranked by how much they matter. Nothing here is advice."
    )
    note.font.size = Pt(12)
    note.font.color.rgb = colour(INK_LIGHT)

    for index, finding in enumerate(_slides(findings), start=1):
        slide = deck.slides.add_slide(blank)
        background(slide)

        eyebrow = slide.shapes.add_textbox(
            Inches(0.9), Inches(0.5), Inches(6), Inches(0.4)
        )
        para = eyebrow.text_frame.paragraphs[0]
        para.text = f"FINDING {index}"
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = colour(_severity_colour(finding.severity))

        heading = slide.shapes.add_textbox(
            Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.3)
        )
        heading.text_frame.word_wrap = True
        head_para = heading.text_frame.paragraphs[0]
        head_para.text = finding.summary
        head_para.font.size = Pt(22)
        head_para.font.bold = True
        head_para.font.color.rgb = colour(INK)

        series = series_for(finding)
        if series and series.values:
            chart_data = CategoryChartData()
            chart_data.categories = series.labels or [
                str(i) for i in range(len(series.values))
            ]
            chart_data.add_series(finding.type.value, series.values)

            kind = {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.DOUGHNUT,
            }[series.kind]

            graphic = slide.shapes.add_chart(
                kind,
                Inches(0.9),
                Inches(2.5),
                Inches(11.5),
                Inches(4.0),
                chart_data,
            )
            chart = graphic.chart
            chart.has_title = False
            if series.kind == "pie":
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.legend.include_in_layout = False
            else:
                chart.has_legend = False
                try:
                    plot = chart.plots[0]
                    plot.series[0].format.fill.solid()
                    plot.series[0].format.fill.fore_color.rgb = colour(ACCENT)
                except (IndexError, AttributeError, NotImplementedError):
                    pass  # a line chart has no fill to set
        else:
            rows = table_for(finding)
            if rows:
                table_shape = slide.shapes.add_table(
                    len(rows), 2, Inches(0.9), Inches(2.6), Inches(7.5),
                    Inches(0.42 * len(rows)),
                )
                table = table_shape.table
                for row_index, (label, value) in enumerate(rows):
                    for column, text in ((0, label), (1, value)):
                        cell = table.cell(row_index, column)
                        cell.text = text
                        para = cell.text_frame.paragraphs[0]
                        para.font.size = Pt(12)
                        para.font.color.rgb = colour(
                            INK_MUTED if column == 0 else INK
                        )
                        if column == 1:
                            para.alignment = PP_ALIGN.RIGHT

        footer = slide.shapes.add_textbox(
            Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4)
        )
        foot = footer.text_frame.paragraphs[0]
        foot.text = f"{finding.evidence.method} · {finding.evidence.strength}"
        foot.font.size = Pt(10)
        foot.font.italic = True
        foot.font.color.rgb = colour(INK_LIGHT)

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()
